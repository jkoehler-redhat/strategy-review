#!/usr/bin/env python3
"""AICP Strategy Review — v2. STRATs-first structure matching leadership review format."""
import sys
sys.path.insert(0, '/opt/homebrew/lib/python3.13/site-packages')

import json, urllib.request, urllib.parse, base64, zipfile, xml.etree.ElementTree as ET
import re, os
from datetime import datetime
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Credentials ──────────────────────────────────────────────────────────────
# Set JIRA_EMAIL / JIRA_TOKEN env vars, or they are read from ~/.claude.json
import os as _os
JIRA_EMAIL = _os.environ.get("JIRA_EMAIL", "")
JIRA_TOKEN = _os.environ.get("JIRA_TOKEN", "")
if not JIRA_EMAIL or not JIRA_TOKEN:
    try:
        import json as _json
        _cfg = _json.load(open(_os.path.expanduser("~/.claude.json")))
        _args = _cfg["mcpServers"]["mcp-atlassian"]["args"]
        JIRA_EMAIL = _args[_args.index("--jira-username") + 1]
        JIRA_TOKEN = _args[_args.index("--jira-token") + 1]
    except Exception as _e:
        raise SystemExit(f"Set JIRA_EMAIL and JIRA_TOKEN env vars, or configure ~/.claude.json: {_e}")
AUTH = base64.b64encode(f"{JIRA_EMAIL}:{JIRA_TOKEN}".encode()).decode()
BASE = "https://redhat.atlassian.net"
TODAY = datetime.now().strftime("%B %d, %Y")
LOOKBACK = 180

# ── Colours ───────────────────────────────────────────────────────────────────
RED   = RGBColor(0xCC, 0x00, 0x00)
DKGRAY= RGBColor(0x33, 0x33, 0x33)
MDGRAY= RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY= RGBColor(0xF2, 0xF2, 0xF2)
BGRED = RGBColor(0xFF, 0xEE, 0xEE)   # very light red for section headers

# ── Jira helpers ──────────────────────────────────────────────────────────────
def jira_search(jql, fields="key,summary,status,priority,labels,issuetype", max_results=100):
    all_issues = []
    next_token = None
    while True:
        params = {"jql": jql, "fields": fields, "maxResults": max_results}
        if next_token: params["nextPageToken"] = next_token
        url = f"{BASE}/rest/api/3/search/jql?{urllib.parse.urlencode(params)}"
        req = urllib.request.Request(url, headers={"Authorization": f"Basic {AUTH}", "Content-Type": "application/json"})
        with urllib.request.urlopen(req) as r: data = json.loads(r.read())
        issues = data.get("issues", [])
        all_issues.extend(issues)
        if data.get("isLast", True) or not issues: break
        next_token = data.get("nextPageToken")
        if not next_token: break
    return all_issues

def jira_get(key, fields="summary,status,priority,description,labels"):
    url = f"{BASE}/rest/api/3/issue/{key}?fields={fields}"
    req = urllib.request.Request(url, headers={"Authorization": f"Basic {AUTH}", "Content-Type": "application/json"})
    with urllib.request.urlopen(req) as r: return json.loads(r.read())["fields"]

def extract_text(node):
    if not node: return ''
    t = node.get('type', ''); content = node.get('content', [])
    if t == 'text':
        text = node.get('text', '')
        for mark in node.get('marks', []):
            if mark.get('type') == 'link':
                href = mark.get('attrs', {}).get('href', '')
                if href and href not in text:
                    text = f"{text} ({href})"
        return text
    if t == 'inlineCard':
        return node.get('attrs', {}).get('url', '')
    if t in ('paragraph', 'heading'): return ' '.join(extract_text(c) for c in content).strip() + '\n'
    if t in ('bulletList', 'orderedList'): return ''.join(extract_text(c) for c in content)
    if t == 'listItem': return '• ' + ' '.join(extract_text(c) for c in content).strip() + '\n'
    if t == 'hardBreak': return '\n'
    return ''.join(extract_text(c) for c in content)

# ── Spreadsheet parser ────────────────────────────────────────────────────────
def parse_roadmap(filepath):
    owned, support = [], []
    with zipfile.ZipFile(filepath) as z:
        shared = []
        if 'xl/sharedStrings.xml' in z.namelist():
            ns = {'s': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
            with z.open('xl/sharedStrings.xml') as f:
                for si in ET.parse(f).getroot().findall('.//s:si', ns):
                    shared.append(''.join(t.text for t in si.findall('.//s:t', ns) if t.text))
        with z.open('xl/workbook.xml') as f:
            wb = ET.parse(f).getroot()
            ns_wb = {'s': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main',
                     'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
            rid = next((sh.get(f'{{{ns_wb["r"]}}}id') for sh in wb.findall('.//s:sheet', ns_wb)
                        if 'Roadmap' in sh.get('name', '')), None)
        sf = None
        with z.open('xl/_rels/workbook.xml.rels') as f:
            for rel in ET.parse(f).getroot().iter():
                if rel.get('Id') == rid: sf = 'xl/' + rel.get('Target'); break
        ns = {'s': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
        with z.open(sf) as f:
            for row in ET.parse(f).getroot().findall('.//s:sheetData/s:row', ns)[1:]:
                rd = {}
                for cell in row.findall('s:c', ns):
                    ref = cell.get('r', ''); col = re.match(r'([A-Z]+)', ref).group(1) if ref else ''
                    t = cell.get('t', ''); v = cell.find('s:v', ns)
                    val = '' if v is None else (shared[int(v.text)] if t == 's' and v.text else (v.text or ''))
                    rd[col] = val
                feat = rd.get('A','').strip(); title = rd.get('B','').strip(); own = rd.get('C','').strip()
                if not feat or not title: continue
                item = {'key': feat, 'title': title, 'ownership': own,
                        'status': rd.get('D','').strip(), 'release': rd.get('F','').strip(),
                        'sub_team': rd.get('J','').strip(), 'priority': rd.get('K','').strip(),
                        'notes': rd.get('M','').strip()}
                # Extract maturity label from title
                mat = 'GA'
                if '[TP]' in title or 'Tech Preview' in title: mat = 'Tech Preview'
                elif '[DP]' in title or 'Dev Preview' in title: mat = 'Dev Preview'
                item['maturity'] = mat
                # Clean title
                item['title_clean'] = re.sub(r'^\[(?:TP|DP)\]\s*', '', title).strip()
                if own.upper() == 'AICP': owned.append(item)
                elif own.upper() == 'SUPPORT': support.append(item)
    return owned, support

# ── Slide helpers ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank_slide(title_text=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Red top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    if title_text:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.65))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title_text
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = DKGRAY; r.font.name = "Calibri"
        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.88), Inches(12), Inches(0.35))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run(); r2.text = subtitle
            r2.font.size = Pt(12); r2.font.color.rgb = MDGRAY; r2.font.name = "Calibri"; r2.font.italic = True
    return slide

def section_slide(title, subtitle=None):
    """Full-bleed section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = DKGRAY; bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    accent.fill.solid(); accent.fill.fore_color.rgb = RED; accent.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(11), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(11), Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); r2.font.name = "Calibri"
    return slide

def add_text_box(slide, text, left, top, width, height, size=13, bold=False, color=None, italic=False, wrap=True):
    color = color or DKGRAY
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_bullet_box(slide, items, left, top, width, height, size=13, space_after=5):
    """items = list of (bold_prefix, rest) tuples."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (bold_part, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        if bold_part:
            r = p.add_run(); r.text = bold_part
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = DKGRAY; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = rest
        r2.font.size = Pt(size); r2.font.color.rgb = DKGRAY; r2.font.name = "Calibri"

def add_table(slide, headers, rows, left, top, width, row_h=0.38, font_size=11):
    n_rows = len(rows) + 1
    shape = slide.shapes.add_table(n_rows, len(headers),
                                   Inches(left), Inches(top), Inches(width), Inches(row_h * n_rows))
    tbl = shape.table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RED
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True
                r.font.size = Pt(font_size); r.font.name = "Calibri"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci); cell.text = str(val)
            if ri % 2 == 1: cell.fill.solid(); cell.fill.fore_color.rgb = LTGRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size); r.font.color.rgb = DKGRAY; r.font.name = "Calibri"
    return tbl

def label_pill(slide, text, left, top, color=RED, text_color=WHITE, width=1.1, height=0.28, size=10):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = text_color; r.font.name = "Calibri"

def bet_slide(title, signal_text, win_text, subtitle=None):
    """Single bet slide with Market Signal + Why This Is Ours to Win."""
    slide = blank_slide(title, subtitle)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); line.line.fill.background()

    # Left column: Market Signal
    add_text_box(slide, "The market signal", 0.6, 1.35, 5.7, 0.4, size=13, bold=True, color=RED)
    add_text_box(slide, signal_text, 0.6, 1.75, 5.7, 5.0, size=12, wrap=True)

    # Vertical divider
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.35), Inches(0.03), Inches(5.4))
    vline.fill.solid(); vline.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); vline.line.fill.background()

    # Right column: Why this is ours to win
    add_text_box(slide, "Why this is ours to win", 6.7, 1.35, 5.9, 0.4, size=13, bold=True, color=RED)
    add_text_box(slide, win_text, 6.7, 1.75, 5.9, 5.0, size=12, wrap=True)
    return slide

# ═══════════════════════════════════════════════════════════════════════════════
# DATA COLLECTION
# ═══════════════════════════════════════════════════════════════════════════════
print("Loading spreadsheet...")
owned, support = parse_roadmap('/private/tmp/strategy-review-repo/aicp_feature_priorities.xlsx')
print(f"  {len(owned)} owned, {len(support)} support")

print("Fetching Jira data...")
inflight = jira_search('project = RHOAIENG AND component in ("AI Core Platform","AI Core Platform Security") AND status in ("In Progress","In Development","In Review","Code Review","In Testing","QE Review")')
blockers  = jira_search('project = RHOAIENG AND component in ("AI Core Platform","AI Core Platform Security") AND priority in (Blocker, Critical) AND status NOT IN (Closed, Resolved, Done, Cancelled)')
upgrade_issues = jira_search('project = RHOAIENG AND component in ("AI Core Platform","AI Core Platform Security") AND summary ~ "upgrade" AND status NOT IN (Closed, Resolved, Done, Cancelled)')

blocker_p = [i for i in blockers if i["fields"].get("priority",{}).get("name") == "Blocker"]
blocker_c = [i for i in blockers if i["fields"].get("priority",{}).get("name") == "Critical"]

def get_team(labels):
    for l in labels:
        if l == "aicp-team-forge": return "Forge"
        if l == "aicp-team-compass": return "Compass"
        if l == "aicp-team-heimdall": return "Heimdall"
    return "Unassigned"

inflight_by_team = defaultdict(list)
for i in inflight: inflight_by_team[get_team(i["fields"].get("labels",[]))].append(i)

# Resolved STRATs (AICP-labeled)
resolved_strats = jira_search(
    f'project = RHAISTRAT AND labels in ("aicp-team-forge","aicp-team-compass","aicp-team-heimdall") AND status in (Closed,Resolved,Done,"Release Pending") AND resolved >= -{LOOKBACK}d',
    fields="key,summary,status,labels")

# Security / TLS compliance epics
security_issues = jira_search(
    'project = RHOAIENG AND component in ("AI Core Platform","AI Core Platform Security") AND (summary ~ "TLS" OR summary ~ "FIPS" OR summary ~ "ML-KEM" OR summary ~ "red-team" OR summary ~ "security scan") AND status NOT IN (Closed, Resolved, Done, Cancelled)',
    fields="key,summary,status,priority,labels")

# Critical bugs (not upgrade, not TLS)
crit_bugs = [i for i in blockers
             if i["fields"].get("priority",{}).get("name") == "Critical"
             and i["fields"].get("issuetype",{}).get("name") in ("Bug","Vulnerability","Weakness")
             and "upgrade" not in i["fields"]["summary"].lower()][:6]

# Component onboarding / near-term deadlines
onboarding = jira_search(
    'project = RHOAIENG AND component in ("AI Core Platform","AI Core Platform Security") AND (summary ~ "operator integration" OR summary ~ "Batch Gateway" OR summary ~ "Agents Operator") AND status NOT IN (Closed, Resolved, Done, Cancelled)',
    fields="key,summary,status,priority,labels")

# EA1 sign-off
ea1 = jira_search(
    'project = RHOAIENG AND component in ("AI Core Platform","AI Core Platform Security") AND summary ~ "Sign-Off" AND status NOT IN (Closed, Resolved, Done, Cancelled)',
    fields="key,summary,status,priority,labels")

# Unlabeled in-flight
unlabeled = [i for i in inflight if get_team(i["fields"].get("labels",[])) == "Unassigned"]

print(f"  In-flight: {len(inflight)} | Blockers: {len(blockers)} | Upgrade: {len(upgrade_issues)} | Security: {len(security_issues)} | Unlabeled: {len(unlabeled)} | Resolved STRATs: {len(resolved_strats)}")

# Fetch strategic bet descriptions
print("Fetching STRAT descriptions for bets...")
def get_desc_snippet(key, max_chars=900):
    try:
        f = jira_get(key, fields="summary,status,description")
        return extract_text(f.get('description') or {})[:max_chars]
    except: return ""

desc_1471 = get_desc_snippet("RHAISTRAT-1471", 900)  # Multi-tenancy
desc_1470 = get_desc_snippet("RHAISTRAT-1470", 900)  # DRA
desc_1519 = get_desc_snippet("RHAISTRAT-1519", 900)  # Upgrade Validation

print("Data collection complete.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
slide = blank_slide()
add_text_box(slide, "AI Core Platform", 0.6, 2.0, 12, 1.0, size=44, bold=True, color=DKGRAY)
add_text_box(slide, "Strategy Review", 0.6, 2.9, 12, 0.9, size=36, color=RED)
add_text_box(slide, f"{TODAY}  |  Red Hat OpenShift AI", 0.6, 3.85, 12, 0.5, size=16, color=MDGRAY)
add_text_box(slide, "Presented by: Jay Koehler, AICP Engineering Management", 0.6, 4.3, 12, 0.4, size=14, color=MDGRAY, italic=True)

# ── SLIDE 2: AGENDA ───────────────────────────────────────────────────────────
slide = blank_slide("Agenda")
add_bullet_box(slide, [
    ("1.  Portfolio Map  ", "— What AICP owns, what we support, and our upgrade investment"),
    ("2.  Traction  ",      "— What we've shipped and what's working"),
    ("3.  Strategic Opportunities  ", "— Three bets backed by customer data, and what we need from leadership"),
    ("",                   ""),
    ("Format:  ",          f"~10 min of slides, then open discussion"),
], left=1.5, top=1.5, width=10, height=4.5, size=17, space_after=12)

# ── SLIDE 3: SECTION — PORTFOLIO MAP ─────────────────────────────────────────
section_slide("Portfolio Map", "What AICP owns, what we support, and where we're investing in platform quality")

# ── SLIDE 4: OWNED STRATs BY TEAM ────────────────────────────────────────────
slide = blank_slide("Owned Roadmap Features", "Portfolio map  ·  Traction  ·  Strategic opportunities")

# Filter to meaningful STRATs (skip internal tracking/Slack/org tickets)
SKIP_KEYS = {"RHOAIENG-41743", "RHOAIENG-54297", "RHOAIENG-23489", "RHOAIENG-9806"}
owned_display = [f for f in owned if f['key'] not in SKIP_KEYS]

# Group by team
by_team = defaultdict(list)
for f in owned_display:
    team = f['sub_team'] if f['sub_team'] in ('Forge','Compass','Heimdall') else 'Platform'
    by_team[team].append(f)

col_width = 3.8; gap = 0.3; start_x = 0.4
for ti, team in enumerate(['Forge', 'Compass', 'Heimdall']):
    x = start_x + ti * (col_width + gap)
    items = by_team.get(team, [])
    # Team header
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.25), Inches(col_width), Inches(0.42))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = RED; hdr.line.fill.background()
    add_text_box(slide, f"{team} ({len(items)})", x+0.08, 1.28, col_width-0.16, 0.38, size=14, bold=True, color=WHITE)
    # Items
    y = 1.75
    for f in items:
        rel = f['release'] if f['release'] and f['release'] != 'N/A' else ''
        mat = f['maturity']
        title = f['title_clean']
        status = f['status']

        # Status dot
        dot_color = RGBColor(0x00,0x99,0x00) if 'Progress' in status else (
                    RGBColor(0xFF,0x99,0x00) if status in ('New','Backlog') else
                    RGBColor(0x00,0x66,0xCC))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.08), Inches(y+0.06), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_color; dot.line.fill.background()

        # Feature title — allow 2 lines of wrap
        add_text_box(slide, title, x+0.25, y, col_width-0.55, 0.52, size=11, color=DKGRAY)

        # Release pill
        if rel:
            add_text_box(slide, rel, x+col_width-0.55, y+0.02, 0.5, 0.22, size=9, color=MDGRAY, bold=True)

        y += 0.58
        if y > 6.8: break

    # For Heimdall: append security program as additional line
    if team == 'Heimdall' and y <= 6.8:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.08), Inches(y+0.06), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(0x00,0x99,0x00); dot.line.fill.background()
        add_text_box(slide, f"TLS / ML-KEM Compliance Program ({len(security_issues)} issues)",
                     x+0.25, y, col_width-0.55, 0.52, size=11, color=DKGRAY)
        add_text_box(slide, "3.5", x+col_width-0.55, y+0.02, 0.5, 0.22, size=9, color=MDGRAY, bold=True)

# Legend
add_text_box(slide, "● In Progress   ● New/Backlog   ● Review/Release Pending", 0.4, 7.05, 12, 0.3, size=9, color=MDGRAY)

# ── SLIDE 5: PLATFORM SUPPORT ENGAGEMENTS ────────────────────────────────────
slide = blank_slide("Platform Support Engagements", "Portfolio map  ·  Traction  ·  Strategic opportunities")

support_display = [f for f in support if f['key'].startswith('RHAISTRAT') or f['key'].startswith('RHAIRFE')]
total_owned = len(owned_display)
total_support = len(support_display)

# Headline stat + framing
add_text_box(slide,
    f"{total_support} support engagements  ·  {total_owned} owned features",
    0.5, 1.25, 9, 0.4, size=18, bold=True, color=DKGRAY)
add_text_box(slide,
    f"Nearly half the AICP portfolio is cross-team support — a capacity and strategic focus tradeoff worth aligning on.",
    0.5, 1.68, 12.3, 0.4, size=12, italic=True, color=MDGRAY)

# Sub-team summary table
by_st = defaultdict(list)
for f in support_display:
    t = f['sub_team'] if f['sub_team'] in ('Forge', 'Compass', 'Heimdall') else 'Shared / Unassigned'
    by_st[t].append(f)

st_rows = []
for team in ['Forge', 'Compass', 'Heimdall', 'Shared / Unassigned']:
    items = by_st.get(team, [])
    if not items: continue
    releases = sorted(set(f['release'] for f in items if f['release'] and f['release'] not in ('N/A', '')))
    rel_str = ', '.join(releases[:3]) if releases else '—'
    notable = next((f['title_clean'] for f in items if f['release'] and f['release'] not in ('N/A', '')),
                   items[0]['title_clean'])
    st_rows.append([team, str(len(items)), rel_str, notable])

add_table(slide, ["Sub-Team", "#", "Releases", "Representative Example"],
          st_rows, left=0.5, top=2.15, width=12.3, row_h=0.42, font_size=11)

# Strategic callouts — pick items with nearest release target
top_support = sorted([f for f in support_display if f['release'] and f['release'] not in ('N/A', '')],
                     key=lambda f: f['release'])[:3]
if len(top_support) < 3:
    top_support += [f for f in support_display if f not in top_support][:3 - len(top_support)]

add_text_box(slide, "Nearest-term commitments", 0.5, 4.42, 5, 0.35, size=12, bold=True, color=DKGRAY)
callouts = []
for f in top_support:
    rel = f['release'] if f['release'] and f['release'] not in ('N/A', '') else 'TBD'
    team = f['sub_team'] if f['sub_team'] else 'Unassigned'
    callouts.append((f"{f['key']}  ", f"{f['title_clean']}  ·  {team}  ·  {rel}"))
add_bullet_box(slide, callouts, left=0.5, top=4.82, width=12.3, height=2.3, size=11, space_after=6)

# ── SLIDE 6: SECTION — TRACTION ───────────────────────────────────────────────
section_slide("Traction", "What's shipped, what's working, and where we're investing")

# ── SLIDE 7: PLATFORM UPGRADE INVESTMENT ──────────────────────────────────────
slide = blank_slide("Platform Upgrade Investment", "Traction  ·  Strategic opportunities")

add_text_box(slide, "RHAISTRAT-1519  ·  Automated Upgrade Validation  ·  Heimdall  ·  3.5  ·  In Progress",
             0.5, 1.25, 12.3, 0.35, size=12, bold=True, color=RED)
add_text_box(slide,
    "Establishes automated upgrade validation as a release quality gate. Every RHOAI release artifact — "
    "nightly, early access, or GA — will execute the supported upgrade matrix, validating both component-owned "
    "and cross-component upgrade scenarios so regressions are caught before they reach customers.",
    0.5, 1.65, 12.3, 0.9, size=12)

# Upgrade engineering work
add_text_box(slide, f"Supporting Engineering  ({len(upgrade_issues)} open RHOAIENG issues)", 0.5, 2.65, 8, 0.35, size=12, bold=True, color=DKGRAY)

up_rows = []
for i in upgrade_issues[:12]:
    s = i["fields"]["summary"]
    for pfx in ("[Upgrade Testing] ", "[RHOAI 3.5.0] ", "[odh-cli 2.25 to 3.3] ", "[QoL Crucible] ", "[Spike] "):
        s = s.replace(pfx, "")
    up_rows.append([i["key"], s, i["fields"]["status"]["name"]])

add_table(slide, ["Issue", "Summary", "Status"],
          up_rows, left=0.5, top=3.05, width=12.3, row_h=0.3, font_size=10)

# ── SLIDE 8: WHAT WE'VE SHIPPED ───────────────────────────────────────────────
slide = blank_slide("What We've Shipped", "Traction  ·  Strategic opportunities")

# Group resolved STRATs thematically
THEME_MAP = {
    "xKS / Multi-Cloud": ["xKS", "kubernetes cluster", "xks", "CoreWeave", "AKS", "distributed inference", "llm-d"],
    "Auth & Gateway":    ["BYOIDC", "OIDC", "gateway", "auth", "Feature Store OIDC", "RBAC", "Granular RBAC"],
    "GitOps":            ["GitOps", "gitops"],
    "Observability":     ["observability", "monitoring", "dashboards", "metrics", "tracing", "Thanos"],
    "MaaS Support":      ["MaaS", "Showback", "External OIDC", "metering"],
    "MLOps":             ["MLflow", "AutoML", "AutoRAG", "Autoscaling", "llm-d autoscal"],
}
def strat_theme(s):
    sl = s.lower()
    for theme, kws in THEME_MAP.items():
        for kw in kws:
            if kw.lower() in sl: return theme
    return "Platform"

by_theme = defaultdict(list)
for i in resolved_strats:
    by_theme[strat_theme(i["fields"]["summary"])].append(i)

add_text_box(slide, f"{len(resolved_strats)} RHAISTRAT items delivered (last {LOOKBACK} days)  ·  {len(inflight)} engineering issues actively in progress",
             0.5, 1.25, 12, 0.35, size=13, bold=True, color=DKGRAY)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.62), Inches(12.3), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); line.line.fill.background()

items = []
for theme in ["xKS / Multi-Cloud", "Auth & Gateway", "GitOps", "Observability", "MaaS Support", "MLOps", "Platform"]:
    iss = by_theme.get(theme, [])
    if not iss: continue
    keys = ", ".join(i["key"] for i in iss[:4])
    more = f" + {len(iss)-4} more" if len(iss) > 4 else ""
    items.append((f"{theme} ({len(iss)})  ", f"{keys}{more}"))

add_bullet_box(slide, items, left=0.5, top=1.7, width=12.3, height=5.3, size=13, space_after=7)

# ── SLIDE 9: SECTION — STRATEGIC OPPORTUNITIES ───────────────────────────────
section_slide("Strategic Opportunities", "Three bets backed by customer data — and what we need from leadership")

# ── SLIDE 10: BET 1 — MULTI-TENANCY ──────────────────────────────────────────
bet_slide(
    "Bet 1: Multi-Tenancy in RHOAI",
    subtitle="RHAISTRAT-1471  ·  Critical  ·  Forge  ·  3.5  ·  In Progress — 0 active engineering issues",
    signal_text=(
        "RHOAI lacks consistent isolation, resource governance, and tenant lifecycle management — forcing "
        "operators to over-provision per tenant or evaluate competing platforms.\n\n"
        "Affected customers:\n"
        "• CCSPs — 20+ Nordic partners (Atea, Advania, Capgemini Finland, CGI, Volvo AB, AI Sweden): "
        "multi-tenancy is a prerequisite for commercializing RHOAI-based AI Factory services.\n"
        "• HPC / Supercomputing — Spain AI Gigafactory, EU AI Continent Action Plan: fair-share GPU "
        "scheduling required by grant funding agencies.\n"
        "• Enterprise — BBVA Garanti, Telenor, Aramco, Nextera, SWIFT: serving 10–50+ internal teams "
        "with independent governance, isolation, and chargeback."
    ),
    win_text=(
        "AICP owns this feature. Multi-tenancy is a layered platform discipline — touching namespaces, "
        "RBAC, networking, resource management, pod security, and cost allocation. This is squarely in "
        "the operator's domain, not any single component.\n\n"
        "Competitive gap: Customers migrating from VMware expect the same ease of dynamic resource "
        "sharing they had with vSphere — virtual clusters, tenant isolation, and on-demand GPU allocation "
        "managed through a single control plane.\n\n"
        "Risk: Feature is Critical priority for 3.5 with zero active engineering work started. "
        "Without resourcing, this becomes a 3.6 slip — losing named accounts."
    )
)

# ── SLIDE 11: BET 2 — GPUaaS / DRA ───────────────────────────────────────────
bet_slide(
    "Bet 2: GPUaaS / Dynamic Resource Allocation (DRA)",
    subtitle="RHAISTRAT-1470  ·  Major  ·  Forge  ·  3.5  ·  In Progress — 2 backlog AICP issues",
    signal_text=(
        "Customers with on-premises NVIDIA GPU investments on OpenShift 4.21 need RHOAI components to "
        "support DRA so GPU resources can be requested with device-level granularity — fractions, specific "
        "attributes, sharing, and multi-device topology.\n\n"
        "Today, RHOAI components (KServe, training-operator, notebooks, KubeRay) use the legacy "
        "nvidia.com/gpu extended resource model, which only supports whole-GPU allocation.\n\n"
        "Affected customers:\n"
        "• 5+ active opportunities in EMEA running OCP 4.21 with centralized AI factory infrastructure\n"
        "• GPUaaS UX vision actively in progress — product expects platform support to land in 3.5\n"
        "• Competitive analysis underway (RHOAIENG-65304)"
    ),
    win_text=(
        "AICP owns the DRA integration work. RHOAI components must be updated to emit DRA ResourceClaims "
        "in pod specs — enabling GPU fractions, sharing, and multi-device topology through the NVIDIA DRA "
        "driver that OCP 4.21 provides.\n\n"
        "This is a platform-layer change: once AICP updates component CRDs and pod specs, every downstream "
        "workload type (training, inference, notebooks) benefits automatically.\n\n"
        "Current state: RHAISTRAT-1470 is In Progress but only 2 AICP engineering issues exist — both in "
        "backlog. The STRAT is moving faster than the engineering work beneath it."
    )
)

# ── SLIDE 12: BET 3 — AUTOMATED UPGRADE VALIDATION ───────────────────────────
bet_slide(
    "Bet 3: Automated Upgrade Validation as a Release Gate",
    subtitle="RHAISTRAT-1519  ·  Heimdall  ·  3.5  ·  In Progress — 9 supporting RHOAIENG issues",
    signal_text=(
        "RHOAI currently lacks a continuous upgrade quality gate in CI. Upgrade regressions are found "
        "late because the migration process still depends on manual steps and limited end-to-end validation.\n\n"
        "This is a recurring source of customer-facing issues at GA — components that pass their own "
        "tests fail during operator upgrades because cross-component upgrade scenarios are not validated "
        "in CI until a release candidate is produced.\n\n"
        f"Supporting engineering work: {len(upgrade_issues)} open RHOAIENG upgrade issues across "
        "Kueue upgrade testing, Gateway API upgrade behavior, xKS upgrade paths, OLM bundle builds, "
        "and disconnected cluster upgrade cases."
    ),
    win_text=(
        "Every produced release artifact — nightly, early access, or GA — will execute the full supported "
        "upgrade matrix and validate both component-owned and cross-component upgrade scenarios.\n\n"
        "For each supported upgrade path, the system will: deploy a pre-upgrade cluster, run pre-upgrade "
        "scenarios, execute the platform upgrade, then validate component and cross-component workloads.\n\n"
        "This directly reduces escalations and shortens time-to-resolution for upgrade-related support "
        "cases. It also unblocks other teams — once AICP owns the upgrade gate, component teams can "
        "own their upgrade validation without manual coordination."
    )
)

# ── SLIDE 13: WHAT WE NEED ────────────────────────────────────────────────────
slide = blank_slide("What We Need", "Portfolio map  ·  Traction  ·  Strategic opportunities")

needs = [
    ("1.  Multi-Tenancy resourcing  ",
     "RHAISTRAT-1471 is Critical for 3.5 with zero active engineering. Need team capacity "
     "assigned to Forge to begin scope definition and delivery this sprint."),
    ("2.  DRA / GPUaaS engineering alignment  ",
     "RHAISTRAT-1470 is In Progress strategically but only 2 backlog issues exist in AICP. "
     "Need engineering investment aligned to product priority — GPUaaS UX vision is already moving."),
    ("3.  Upgrade CI infrastructure  ",
     f"RHAISTRAT-1519 (Heimdall) requires cluster provisioning infrastructure to run the upgrade "
     f"matrix in CI. {len(upgrade_issues)} supporting issues in flight — need platform/infra support."),
    ("4.  Blocker escalation  ",
     f"{len(blocker_p)} blocker-priority and {len(blocker_c)} critical issues remain open across AICP. "
     "Top: xKS operator CrashLoopBackOff (RHOAIENG-67048), LLM-D Batch Gateway integration "
     "(RHOAIENG-63970/63964), FIPS compliance (RHOAIENG-63211)."),
    ("5.  Sub-team triage  ",
     f"{len(unlabeled)} in-flight issues — including {len([i for i in security_issues if get_team(i['fields'].get('labels',[])) == 'Unassigned'])} of the TLS compliance epics — have no sub-team label. "
     "Need manager review to assign to Forge, Compass, or Heimdall so they appear in reporting and sprint planning."),
]
add_bullet_box(slide, needs, left=0.5, top=1.3, width=12.3, height=5.8, size=13, space_after=8)

# ── SLIDE 14: THANK YOU ───────────────────────────────────────────────────────
slide = blank_slide()
add_text_box(slide, "Thank you", 0.6, 2.5, 12, 1.0, size=44, bold=True, color=DKGRAY)
add_text_box(slide, "Open discussion", 0.6, 3.5, 12, 0.7, size=28, color=RED)
add_text_box(slide, "What questions do you have?  Where should we go deeper?", 0.6, 4.3, 12, 0.5, size=16, color=MDGRAY, italic=True)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out_dir = "/Users/jaykoehler/aicp-status/docs/strategy"
os.makedirs(out_dir, exist_ok=True)
date_str = datetime.now().strftime("%Y-%m-%d")
path1 = f"{out_dir}/{date_str}_strategy_review.pptx"
path2 = "/Users/jaykoehler/aicp-strategy-review.pptx"
prs.save(path1); print(f"Saved: {path1}")
prs.save(path2); print(f"Saved: {path2}")
print(f"\n14 slides: Title · Agenda · [Portfolio Map] · Owned STRATs · Support · [Traction] · Upgrade · Shipped · [Opportunities] · Bet1 · Bet2 · Bet3 · What We Need · Thank You")
