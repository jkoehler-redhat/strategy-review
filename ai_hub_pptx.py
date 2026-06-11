#!/usr/bin/env python3
"""AI Hub Strategy Review — same structure as AICP deck, scoped to component = "AI Hub"."""
import sys
sys.path.insert(0, '/opt/homebrew/lib/python3.13/site-packages')

import json, urllib.request, urllib.parse, base64, re, os
from datetime import datetime
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Credentials ───────────────────────────────────────────────────────────────
import os as _os
JIRA_EMAIL = _os.environ.get("JIRA_EMAIL", "")
JIRA_TOKEN = _os.environ.get("JIRA_TOKEN", "")
if not JIRA_EMAIL or not JIRA_TOKEN:
    try:
        import json as _json
        _cfg = _json.load(open(_os.path.expanduser("~/.claude.json")))
        _args = _cfg["mcpServers"]["mcp-atlassian"]["args"]
        JIRA_EMAIL = _args[_args.index("--jira-username") + 1]
        JIRA_TOKEN = _args[_args.index("--jira-token") + 1]
    except Exception as _e:
        raise SystemExit(f"Set JIRA_EMAIL and JIRA_TOKEN env vars, or configure ~/.claude.json: {_e}")
AUTH = base64.b64encode(f"{JIRA_EMAIL}:{JIRA_TOKEN}".encode()).decode()
BASE = "https://redhat.atlassian.net"
TODAY = datetime.now().strftime("%B %d, %Y")
LOOKBACK = 180

# ── Colours ───────────────────────────────────────────────────────────────────
RED   = RGBColor(0xCC, 0x00, 0x00)
DKGRAY= RGBColor(0x33, 0x33, 0x33)
MDGRAY= RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY= RGBColor(0xF2, 0xF2, 0xF2)

# ── Jira helpers ──────────────────────────────────────────────────────────────
def jira_search(jql, fields="key,summary,status,priority,labels,issuetype", max_results=100):
    all_issues = []
    next_token = None
    while True:
        params = {"jql": jql, "fields": fields, "maxResults": max_results}
        if next_token: params["nextPageToken"] = next_token
        url = f"{BASE}/rest/api/3/search/jql?{urllib.parse.urlencode(params)}"
        req = urllib.request.Request(url, headers={"Authorization": f"Basic {AUTH}", "Content-Type": "application/json"})
        with urllib.request.urlopen(req) as r: data = json.loads(r.read())
        issues = data.get("issues", [])
        all_issues.extend(issues)
        if data.get("isLast", True) or not issues: break
        next_token = data.get("nextPageToken")
        if not next_token: break
    return all_issues

def jira_get(key, fields="summary,status,priority,description,labels"):
    url = f"{BASE}/rest/api/3/issue/{key}?fields={fields}"
    req = urllib.request.Request(url, headers={"Authorization": f"Basic {AUTH}", "Content-Type": "application/json"})
    with urllib.request.urlopen(req) as r: return json.loads(r.read())["fields"]

def extract_text(node):
    if not node: return ''
    t = node.get('type', ''); content = node.get('content', [])
    if t == 'text':
        text = node.get('text', '')
        for mark in node.get('marks', []):
            if mark.get('type') == 'link':
                href = mark.get('attrs', {}).get('href', '')
                if href and href not in text:
                    text = f"{text} ({href})"
        return text
    if t == 'inlineCard':
        return node.get('attrs', {}).get('url', '')
    if t in ('paragraph', 'heading'): return ' '.join(extract_text(c) for c in content).strip() + '\n'
    if t in ('bulletList', 'orderedList'): return ''.join(extract_text(c) for c in content)
    if t == 'listItem': return '• ' + ' '.join(extract_text(c) for c in content).strip() + '\n'
    if t == 'hardBreak': return '\n'
    return ''.join(extract_text(c) for c in content)

def get_desc_snippet(key, max_chars=900):
    try:
        f = jira_get(key, fields="summary,status,description")
        return extract_text(f.get('description') or {})[:max_chars]
    except: return ""

# ── Theme mapping ─────────────────────────────────────────────────────────────
THEME_MAP = {
    "Model Catalog":          ["catalog", "eagle3", "xeon", "tool calling", "text-embedding", "cold-start", "vram", "benchmark", "model batch"],
    "MCP Registry":           ["mcp registry", "mcp server"],
    "MLflow / Asset Registry":["mlflow", "unity catalog", "asset registry", "tiger team", "ai asset"],
    "Bodies of Water":        ["bodies of water", "lake", "ocean", "stream", "modular upgrade"],
    "Security & CVEs":        ["cve", "vulnerability", "schemathesis", "ssrf", "hermetic", "fuzz", "signing", "cosign", "securesign"],
    "Async Upload / OCI":     ["async upload", "async-upload", "omlmd", "oci", "image sign", "omlmd"],
    "Upstream / Community":   ["kubeflow", "upstream", "mlmd", "graduation", "community", "blog"],
}
def get_theme(summary):
    sl = summary.lower()
    for theme, kws in THEME_MAP.items():
        for kw in kws:
            if kw.lower() in sl: return theme
    return "Model Registry Core"

# ── Slide helpers ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank_slide(title_text=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    if title_text:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.65))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title_text
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = DKGRAY; r.font.name = "Calibri"
        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.88), Inches(12), Inches(0.35))
            p2 = tb2.text_frame.paragraphs[0]
            r2 = p2.add_run(); r2.text = subtitle
            r2.font.size = Pt(12); r2.font.color.rgb = MDGRAY; r2.font.name = "Calibri"; r2.font.italic = True
    return slide

def section_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = DKGRAY; bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    accent.fill.solid(); accent.fill.fore_color.rgb = RED; accent.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(11), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(11), Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0xCC,0xCC,0xCC); r2.font.name = "Calibri"
    return slide

def add_text_box(slide, text, left, top, width, height, size=13, bold=False, color=None, italic=False, wrap=True):
    color = color or DKGRAY
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def add_bullet_box(slide, items, left, top, width, height, size=13, space_after=5):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (bold_part, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        if bold_part:
            r = p.add_run(); r.text = bold_part
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = DKGRAY; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = rest
        r2.font.size = Pt(size); r2.font.color.rgb = DKGRAY; r2.font.name = "Calibri"

def add_table(slide, headers, rows, left, top, width, row_h=0.38, font_size=11):
    n_rows = len(rows) + 1
    shape = slide.shapes.add_table(n_rows, len(headers),
                                   Inches(left), Inches(top), Inches(width), Inches(row_h * n_rows))
    tbl = shape.table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RED
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True
                r.font.size = Pt(font_size); r.font.name = "Calibri"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci); cell.text = str(val)
            if ri % 2 == 1: cell.fill.solid(); cell.fill.fore_color.rgb = LTGRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size); r.font.color.rgb = DKGRAY; r.font.name = "Calibri"

def bet_slide(title, signal_text, win_text, subtitle=None):
    slide = blank_slide(title, subtitle)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); line.line.fill.background()
    add_text_box(slide, "The market signal", 0.6, 1.35, 5.7, 0.4, size=13, bold=True, color=RED)
    add_text_box(slide, signal_text, 0.6, 1.75, 5.7, 5.0, size=12, wrap=True)
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.35), Inches(0.03), Inches(5.4))
    vline.fill.solid(); vline.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); vline.line.fill.background()
    add_text_box(slide, "Why this is ours to win", 6.7, 1.35, 5.9, 0.4, size=13, bold=True, color=RED)
    add_text_box(slide, win_text, 6.7, 1.75, 5.9, 5.0, size=12, wrap=True)
    return slide

# ═══════════════════════════════════════════════════════════════════════════════
# DATA COLLECTION
# ═══════════════════════════════════════════════════════════════════════════════
print("Fetching AI Hub Jira data...")

COMP = '"AI Hub"'

all_open = jira_search(f'project = RHOAIENG AND component = {COMP} AND status NOT IN (Closed, Resolved, Done, Cancelled)')
inflight  = [i for i in all_open if i["fields"]["status"]["name"] in
             ("In Progress","In Development","In Review","Code Review","In Testing","QE Review","Review","Testing")]
blockers  = [i for i in all_open if i["fields"].get("priority",{}).get("name") in ("Blocker","Critical")]
blocker_p = [i for i in blockers if i["fields"].get("priority",{}).get("name") == "Blocker"]
blocker_c = [i for i in blockers if i["fields"].get("priority",{}).get("name") == "Critical"]
cves      = [i for i in all_open if i["fields"].get("issuetype",{}).get("name") in ("Vulnerability","Weakness")]
undefined = [i for i in all_open if i["fields"].get("priority",{}).get("name") == "Undefined"]

# Group in-flight by theme
inflight_by_theme = defaultdict(list)
for i in inflight:
    inflight_by_theme[get_theme(i["fields"]["summary"])].append(i)

# Resolved in lookback
resolved = jira_search(
    f'project = RHOAIENG AND component = {COMP} AND status in (Closed,Resolved,Done,"Release Pending") AND resolved >= -{LOOKBACK}d',
    fields="key,summary,status,labels,issuetype")

resolved_by_theme = defaultdict(list)
for i in resolved:
    resolved_by_theme[get_theme(i["fields"]["summary"])].append(i)

print(f"  Open: {len(all_open)} | In-flight: {len(inflight)} | Blockers: {len(blocker_p)} | CVEs: {len(cves)} | Resolved ({LOOKBACK}d): {len(resolved)}")

# Strategic bet descriptions
print("Fetching strategic issue descriptions...")
desc_mlflow  = get_desc_snippet("RHOAIENG-50747", 900)  # MLflow tiger team
desc_mcp     = get_desc_snippet("RHOAIENG-63382", 900)  # MCP Registry
desc_catalog = get_desc_snippet("RHOAIENG-60367", 900)  # Decouple catalog

print("Data collection complete.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
slide = blank_slide()
add_text_box(slide, "AI Hub", 0.6, 2.0, 12, 1.0, size=44, bold=True, color=DKGRAY)
add_text_box(slide, "Strategy Review", 0.6, 2.9, 12, 0.9, size=36, color=RED)
add_text_box(slide, f"{TODAY}  |  Red Hat OpenShift AI", 0.6, 3.85, 12, 0.5, size=16, color=MDGRAY)
add_text_box(slide, "Presented by: Jay Koehler, AICP Engineering Management", 0.6, 4.3, 12, 0.4, size=14, color=MDGRAY, italic=True)

# ── SLIDE 2: AGENDA ───────────────────────────────────────────────────────────
slide = blank_slide("Agenda")
add_bullet_box(slide, [
    ("1.  Portfolio Map  ", "— What AI Hub owns and what's in flight"),
    ("2.  Traction  ",      "— What we've shipped"),
    ("3.  Strategic Opportunities  ", "— Three bets backed by data, and what we need from leadership"),
    ("",                   ""),
    ("Format:  ",          "~10 min of slides, then open discussion"),
], left=1.5, top=1.5, width=10, height=4.5, size=17, space_after=12)

# ── SLIDE 3: SECTION — PORTFOLIO MAP ─────────────────────────────────────────
section_slide("Portfolio Map", "What AI Hub owns and what's actively in flight")

# ── SLIDE 4: PORTFOLIO OVERVIEW ───────────────────────────────────────────────
slide = blank_slide("Portfolio Overview", "Portfolio map  ·  Traction  ·  Strategic opportunities")

add_text_box(slide,
    f"{len(all_open)} open issues  ·  {len(inflight)} actively in progress  ·  {len(blocker_p)} blockers  ·  {len(blocker_c)} critical  ·  {len(cves)} open CVEs",
    0.5, 1.25, 12.3, 0.4, size=14, bold=True, color=DKGRAY)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.72), Inches(12.3), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); line.line.fill.background()

# Theme breakdown table
theme_rows = []
all_themes = ["Model Catalog", "MCP Registry", "MLflow / Asset Registry", "Bodies of Water",
              "Security & CVEs", "Async Upload / OCI", "Upstream / Community", "Model Registry Core"]
for theme in all_themes:
    open_count = sum(1 for i in all_open if get_theme(i["fields"]["summary"]) == theme)
    if open_count == 0: continue
    flight_count = len(inflight_by_theme.get(theme, []))
    bloc_count = sum(1 for i in blockers if get_theme(i["fields"]["summary"]) == theme)
    res_count = len(resolved_by_theme.get(theme, []))
    theme_rows.append([theme, str(open_count), str(flight_count), str(bloc_count), str(res_count)])

add_table(slide, ["Theme", "Open", "In Flight", "Blockers/Critical", f"Resolved ({LOOKBACK}d)"],
          theme_rows, left=0.5, top=1.82, width=12.3, row_h=0.42, font_size=11)

add_text_box(slide, f"{len(undefined)} issues have Undefined priority — need triage",
             0.5, 7.05, 12.3, 0.3, size=10, color=MDGRAY, italic=True)

# ── SLIDE 5: IN FLIGHT BY THEME ───────────────────────────────────────────────
slide = blank_slide("What's In Flight", "Portfolio map  ·  Traction  ·  Strategic opportunities")

add_text_box(slide, f"{len(inflight)} issues actively in progress across {len([t for t in inflight_by_theme if inflight_by_theme[t]])} themes",
             0.5, 1.25, 12.3, 0.35, size=13, bold=True, color=DKGRAY)

items = []
for theme in all_themes:
    iss = inflight_by_theme.get(theme, [])
    if not iss: continue
    top = iss[:3]
    bullets = "; ".join(i["fields"]["summary"] for i in top)
    more = f" + {len(iss)-3} more" if len(iss) > 3 else ""
    items.append((f"{theme} ({len(iss)})  ", f"{bullets}{more}"))

add_bullet_box(slide, items, left=0.5, top=1.7, width=12.3, height=5.3, size=12, space_after=8)

# ── SLIDE 6: SECTION — TRACTION ───────────────────────────────────────────────
section_slide("Traction", "What's shipped in the last 6 months")

# ── SLIDE 7: WHAT WE'VE SHIPPED ───────────────────────────────────────────────
slide = blank_slide("What We've Shipped", "Traction  ·  Strategic opportunities")

add_text_box(slide, f"{len(resolved)} issues resolved in the last {LOOKBACK} days",
             0.5, 1.25, 12.3, 0.35, size=13, bold=True, color=DKGRAY)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); line.line.fill.background()

items = []
for theme in all_themes:
    iss = resolved_by_theme.get(theme, [])
    if not iss: continue
    keys = ", ".join(i["key"] for i in iss[:4])
    more = f" + {len(iss)-4} more" if len(iss) > 4 else ""
    items.append((f"{theme} ({len(iss)})  ", f"{keys}{more}"))

add_bullet_box(slide, items, left=0.5, top=1.75, width=12.3, height=5.3, size=13, space_after=8)

# ── SLIDE 8: SECTION — STRATEGIC OPPORTUNITIES ───────────────────────────────
section_slide("Strategic Opportunities", "Three bets backed by data — and what we need from leadership")

# ── SLIDE 9: BET 1 — MLFLOW / AI ASSET REGISTRY ─────────────────────────────
bet_slide(
    "Bet 1: MLflow as the Unified AI Asset Registry",
    subtitle="RHOAIENG-50747  ·  Tiger Team  ·  In Progress",
    signal_text=(
        "Kubeflow Model Registry is a narrowly scoped model metadata store. Customers need a unified "
        "registry that covers models, experiments, datasets, and artifacts — with lineage, governance, "
        "and multi-tenancy built in.\n\n"
        "MLflow / Unity Catalog OSS has emerged as the leading open-source candidate for this role. "
        "Multiple enterprise customers are already running MLflow and expect RHOAI to integrate with it "
        "rather than requiring migration to a proprietary registry.\n\n"
        + (desc_mlflow if desc_mlflow else
           "Tiger team is actively investigating Unity Catalog as the governance backbone for a unified AI asset registry.")
    ),
    win_text=(
        "AI Hub owns Model Registry — and owns the migration path. Moving to MLflow as the backend "
        "means the team leads the architectural decision that shapes how RHOAI manages AI assets for "
        "the next 3+ years.\n\n"
        "Current state: Tiger team research complete (RHOAIENG-50747). MLflow tracking store PR rebased "
        "upstream (RHOAIENG-31396). RFC in progress (RHOAIENG-57652). Next step: committing to a "
        "migration architecture and engineering plan.\n\n"
        "Risk: Without a decision, teams continue building on top of Kubeflow Model Registry — "
        "accruing integration debt that becomes harder to migrate later."
    )
)

# ── SLIDE 10: BET 2 — MCP REGISTRY ───────────────────────────────────────────
bet_slide(
    "Bet 2: MCP Registry — Model Context Protocol as a Platform Capability",
    subtitle="RHOAIENG-63382 / RHOAIENG-63350  ·  Tech Preview  ·  In Progress",
    signal_text=(
        "Model Context Protocol (MCP) is rapidly becoming the standard for how AI agents discover and "
        "invoke tools, APIs, and data sources. Customers building agentic AI workflows on RHOAI expect "
        "a platform-managed registry for MCP servers — not ad hoc configuration.\n\n"
        "AI Hub is building this capability now: MCP server registration, versioning, RBAC, and Python "
        "SDK integration. Tech Preview is in progress.\n\n"
        + (desc_mcp if desc_mcp else
           "Backend work actively in progress: REST client, Python SDK methods, RBAC for MCP routes (RHOAIENG-63418–63421).")
    ),
    win_text=(
        "AI Hub is positioned to own MCP Registry as a first-class platform service — the same way "
        "it owns Model Registry. This creates a natural extension of the Hub: a single place to discover "
        "models, artifacts, and now agentic tools.\n\n"
        "Current state: Backend tasks in progress (RHOAIENG-63418 REST client in Review, RHOAIENG-63419 "
        "SDK in progress). RBAC and route permissions scoped (RHOAIENG-63421). Tech Preview targeted "
        "for 3.5.\n\n"
        "Opportunity: If AI Hub ships MCP Registry at Tech Preview, it becomes the reference "
        "implementation for how RHOAI manages agentic infrastructure."
    )
)

# ── SLIDE 11: BET 3 — MODEL CATALOG DECOUPLING ───────────────────────────────
bet_slide(
    "Bet 3: Decouple Model Catalog into an Independent Operator",
    subtitle="RHOAIENG-60367  ·  Major  ·  New — stories scoped, not yet started",
    signal_text=(
        "Model Catalog is currently embedded inside the model-registry-operator, creating tight coupling "
        "that limits independent releases, increases blast radius for catalog changes, and complicates "
        "the onboarding of new catalog sources (Hugging Face, OCI, partner catalogs).\n\n"
        "Customers want to deploy and update the model catalog independently of the model registry. "
        "As catalog integrations grow — Eagle3, Xeon, text-embedding, Tool Calling defaults — the "
        "release cadence mismatch becomes a bottleneck.\n\n"
        + (desc_catalog if desc_catalog else
           "Stories scoped: bootstrap catalog-operator (RHOAIENG-60372), port catalog controller (RHOAIENG-60379), "
           "add catalog DSC component (RHOAIENG-60378), deprecate catalog from model-registry-operator (RHOAIENG-60383).")
    ),
    win_text=(
        "Decoupling the catalog gives AI Hub independent release velocity for the fastest-moving part "
        "of the product — new model integrations, catalog sources, and discovery features — without "
        "waiting on model registry release cycles.\n\n"
        "The architecture is already designed: catalog-operator with its own CRD, bootstrapped with "
        "kubebuilder scaffold (RHOAIENG-60372). Four stories fully scoped and ready to start.\n\n"
        "This also sets the foundation for the MLflow migration — a decoupled catalog can adopt "
        "a new registry backend without a full platform change."
    )
)

# ── SLIDE 12: WHAT WE NEED ────────────────────────────────────────────────────
slide = blank_slide("What We Need", "Portfolio map  ·  Traction  ·  Strategic opportunities")

needs = [
    ("1.  MLflow architecture decision  ",
     "Tiger team research is complete. Need leadership alignment to commit to MLflow / Unity Catalog "
     "as the registry backend — or explicitly decide to stay on Kubeflow MR. Delay accrues integration debt."),
    ("2.  MCP Registry Tech Preview resourcing  ",
     "RHOAIENG-63382 backend work is in progress but scoped narrowly. Need dedicated capacity to "
     "complete RBAC, SDK, and QE requirements before the 3.5 TP window closes."),
    ("3.  CVE & security backlog  ",
     f"{len(cves)} open CVEs and vulnerabilities — including embargoed issues. Async upload job hermetic "
     "build failures are blocking several CVE fixes. Need infra support to unblock hermetic builds."),
    ("4.  Priority triage  ",
     f"{len(undefined)} issues ({round(len(undefined)/len(all_open)*100)}% of portfolio) have Undefined priority. "
     "Reporting, sprint planning, and escalation paths are all unreliable until these are triaged."),
]
add_bullet_box(slide, needs, left=0.5, top=1.3, width=12.3, height=5.8, size=13, space_after=10)

# ── SLIDE 13: THANK YOU ───────────────────────────────────────────────────────
slide = blank_slide()
add_text_box(slide, "Thank you", 0.6, 2.5, 12, 1.0, size=44, bold=True, color=DKGRAY)
add_text_box(slide, "Open discussion", 0.6, 3.5, 12, 0.7, size=28, color=RED)
add_text_box(slide, "What questions do you have?  Where should we go deeper?", 0.6, 4.3, 12, 0.5, size=16, color=MDGRAY, italic=True)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out_dir = "/Users/jaykoehler/aicp-status/docs/strategy"
os.makedirs(out_dir, exist_ok=True)
date_str = datetime.now().strftime("%Y-%m-%d")
path1 = f"{out_dir}/{date_str}_ai_hub_strategy_review.pptx"
path2 = "/Users/jaykoehler/ai-hub-strategy-review.pptx"
prs.save(path1); print(f"Saved: {path1}")
prs.save(path2); print(f"Saved: {path2}")
print(f"\n13 slides: Title · Agenda · [Portfolio Map] · Overview · In Flight · [Traction] · Shipped · [Opportunities] · Bet1 · Bet2 · Bet3 · What We Need · Thank You")
